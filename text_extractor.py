"""
text_extractor.py - Document Text Extraction
=============================================
Handles text extraction from PDF, DOCX, PPTX, and image files.
Images are processed via Tesseract OCR.
"""

import os
import re
import logging

logger = logging.getLogger(__name__)


# ── Helpers ───────────────────────────────────────────────────────────────────

def _clean_text(text: str) -> str:
    """
    Normalise extracted text:
    - Collapse excessive whitespace / blank lines
    - Strip leading/trailing whitespace per line
    """
    lines = [line.strip() for line in text.splitlines()]
    # Remove sequences of more than two consecutive blank lines
    cleaned_lines: list[str] = []
    blank_streak = 0
    for line in lines:
        if line == "":
            blank_streak += 1
            if blank_streak <= 2:
                cleaned_lines.append(line)
        else:
            blank_streak = 0
            cleaned_lines.append(line)
    return "\n".join(cleaned_lines).strip()


# ── Extractors ────────────────────────────────────────────────────────────────

def extract_from_pdf(filepath: str) -> str:
    """Extract text from a PDF file using PyMuPDF (fitz)."""
    try:
        import fitz  # PyMuPDF
        text_parts: list[str] = []
        with fitz.open(filepath) as doc:
            for page in doc:
                text_parts.append(page.get_text("text"))
        return _clean_text("\n".join(text_parts))
    except Exception as exc:
        logger.error("PDF extraction failed: %s", exc)
        raise RuntimeError(f"Could not extract text from PDF: {exc}") from exc


def extract_from_docx(filepath: str) -> str:
    """Extract text from a DOCX file using python-docx."""
    try:
        from docx import Document
        doc = Document(filepath)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also pull text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        paragraphs.append(cell.text.strip())
        return _clean_text("\n".join(paragraphs))
    except Exception as exc:
        logger.error("DOCX extraction failed: %s", exc)
        raise RuntimeError(f"Could not extract text from DOCX: {exc}") from exc


def extract_from_pptx(filepath: str) -> str:
    """Extract text from a PPTX file using python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        text_parts: list[str] = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            text_parts.append(f"--- Slide {slide_num} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs).strip()
                        if line:
                            text_parts.append(line)
        return _clean_text("\n".join(text_parts))
    except Exception as exc:
        logger.error("PPTX extraction failed: %s", exc)
        raise RuntimeError(f"Could not extract text from PPTX: {exc}") from exc


def extract_from_image(filepath: str) -> str:
    """
    Run Tesseract OCR on an image file.
    Requires Tesseract to be installed on the host system.
    """
    try:
        import pytesseract
        from PIL import Image

        img = Image.open(filepath)
        # Convert to RGB to avoid mode issues (e.g. RGBA PNGs)
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        text = pytesseract.image_to_string(img, lang="eng")
        return _clean_text(text)
    except Exception as exc:
        logger.error("OCR extraction failed: %s", exc)
        raise RuntimeError(f"Could not extract text via OCR: {exc}") from exc


# ── Public API ────────────────────────────────────────────────────────────────

def extract_text(filepath: str, extension: str) -> str:
    """
    Dispatch to the correct extractor based on file extension.

    Args:
        filepath:  Absolute path to the uploaded file.
        extension: Lowercase extension without the leading dot (e.g. "pdf").

    Returns:
        Cleaned extracted text string.

    Raises:
        ValueError: If the extension is unsupported.
        RuntimeError: If extraction fails.
    """
    ext = extension.lower().lstrip(".")

    extractors = {
        "pdf":  extract_from_pdf,
        "docx": extract_from_docx,
        "pptx": extract_from_pptx,
        "png":  extract_from_image,
        "jpg":  extract_from_image,
        "jpeg": extract_from_image,
    }

    if ext not in extractors:
        raise ValueError(f"Unsupported file type: .{ext}")

    text = extractors[ext](filepath)

    if not text or len(text.strip()) < 20:
        raise RuntimeError(
            "Extracted text is too short or empty. "
            "Please upload a document with readable content."
        )

    return text
